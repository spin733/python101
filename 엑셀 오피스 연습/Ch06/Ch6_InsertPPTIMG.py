from pptx import Presentation
from tkinter import filedialog
import os

#프리젠 테이션 열기
prs = Presentation('Ch6_Sample.pptx')

#폴더 내 파일리스트 가져오기
foldername = filedialog.askdirectory()
filearr = os.listdir(foldername)

#파일 반복문 실행
for file in filearr:
    slide = prs.slides.add_slide(prs.slide_layouts[2])
    slide.placeholders[10].insert_picture(foldername + '/' + file)
    slide.placeholders[11].text = file

prs.save('Ch6_Result.pptx')