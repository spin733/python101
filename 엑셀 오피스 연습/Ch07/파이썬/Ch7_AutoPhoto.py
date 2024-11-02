import pandas as pd

df = pd.read_excel('Ch7_SpeedTest.xlsx')

from pptx import Presentation
import os

prs = Presentation('Ch7_AutoPhoto_Format.pptx')

slide = prs.slides.add_slide(prs.slide_layouts[2])

for i in df.index:
	odev = (df.No[i] +1) % 2
	filename = os.path.join(df.Path[i], df.FileName[i])
	slide.placeholders[10 + odev * 5].insert_picture(filename)
	slide.placeholders[11 + odev * 5].text = df.ConName[i]
	slide.placeholders[12 + odev * 5].text = df.Description[i]
	slide.placeholders[13 + odev * 5].text = df.Place[i]
	slide.placeholders[14 + odev * 5].text = df.Date[i].strftime('%Y/%m/%d')
	try:
		odev = (df.No[i+1] + 1) % 2
		if odev == 0:
			slide = prs.slides.add_slide(prs.slide_layouts[2])
	except:
		pass
prs.save('Ch7_Result.pptx')
