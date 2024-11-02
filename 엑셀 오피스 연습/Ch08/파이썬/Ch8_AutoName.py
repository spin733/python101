import pandas as pd

df = pd.read_excel('Ch8_NameCardList.xlsx')

from pptx import Presentation
import os

prs = Presentation('Ch8_AutoNameCard.pptx')

from pptx.enum.text import PP_ALIGN

# paragraphs[0].alignment = PP_ALIGN.CENTER

for i in df.index:
	odev = i % 8
	filename = os.path.join("symbol/", df.Department[i] + ".png")

	if odev == 0 :
		slide = prs.slides.add_slide(prs.slide_layouts[2])

	slide.placeholders[13 + odev * 4].insert_picture(filename)
		# 10 소속 Department
		# 11 이름 Name
		# 12 직책 Position
		# 13 Symbol "Symbol/" Department ".png"
	slide.placeholders[10 + odev * 4].text = df.Name[i]
	slide.placeholders[10 + odev * 4].text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
	slide.placeholders[11 + odev * 4].text = df.Department[i]
	slide.placeholders[12 + odev * 4].text = df.Position[i]
	slide.placeholders[12 + odev * 4].text_frame.paragraphs[0].alignment = PP_ALIGN.RIGHT

prs.save('Ch8_NameResult.pptx')


