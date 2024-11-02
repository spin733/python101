import pandas as pd

df = pd.read_excel('Ch8_AwardList.xlsx')

from pptx import Presentation
import os

prs = Presentation('Ch8_AutoAward_Format.pptx')

import datetime as dt

df['real_date'] = pd.TimedeltaIndex(df.Date, unit='d') + dt.datetime(1899, 12, 30)

for i in df.index:
	#odev = (df.No[i] +1) % 2
	filename = os.path.join("symbol/", df.Department[i] + ".png")

	slide = prs.slides.add_slide(prs.slide_layouts[1])
	slide.placeholders[15].insert_picture(filename)
	# 10 : No
	# 11 : Part
	# 12 : Grade
	# 13 : Department
	# 14 : Name
	# 16 : Date
	slide.placeholders[10].text = str(df.No[i])
	slide.placeholders[11].text = df.Part[i]
	slide.placeholders[12].text = df.Grade[i]
	slide.placeholders[13].text = df.Department[i]
	slide.placeholders[14].text = df.Name[i]
	slide.placeholders[16].text = df.real_date[i].strftime('%Y/%m/%d')
	# try:
	# 	odev = (df.No[i+1] + 1) % 2
	# 	if odev == 0:
	#
	# except:
	# 	pass
prs.save('Ch8_Result.pptx')

